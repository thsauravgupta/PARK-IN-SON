import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    NAVY = RGBColor(15, 23, 42)       # Slate 900
    LIGHT_GRAY = RGBColor(248, 250, 252) # Slate 50
    TEXT_DARK = RGBColor(30, 41, 59)   # Slate 800
    TEXT_LIGHT = RGBColor(241, 245, 249) # Slate 100
    GOLD = RGBColor(217, 119, 6)       # Amber 600
    TEAL = RGBColor(13, 148, 136)      # Teal 600
    CARD_BG = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(226, 232, 240) # Slate 200

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title(slide, text, color=TEXT_DARK, size=28, top=0.5):
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.833), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Trebuchet MS"
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color
        return txBox

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        # remove margins
        shape.text_frame.margin_left = Inches(0.2)
        shape.text_frame.margin_top = Inches(0.2)
        shape.text_frame.margin_right = Inches(0.2)
        shape.text_frame.margin_bottom = Inches(0.2)
        return shape

    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Title Slide (Dark Navy)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, NAVY)

    # Title & Subtitle in single textbox to prevent overlap
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0

    p_title = tf.paragraphs[0]
    p_title.text = "Fed-CM-CADT Network"
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = GOLD
    p_title.space_after = Pt(10)

    p_sub = tf.add_paragraph()
    p_sub.text = "Federated Cross-Modal Co-Attentive Dual-Twin Network for Parkinson's Disease Prediction"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = TEXT_LIGHT
    p_sub.space_after = Pt(30)

    p_meta = tf.add_paragraph()
    p_meta.text = "A Publication-Grade Architecture Proposal for IEEE/Elsevier Journals\nMulti-Modal Fusion (Clinical, Genetics, MRI, PET) + Federated Learning + XAI"
    p_meta.font.name = "Arial"
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = RGBColor(148, 163, 184) # Slate 400

    # =========================================================================
    # SLIDE 2: Project Overview & Current Baseline
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, LIGHT_GRAY)
    add_title(slide2, "Phase 1 Baseline & Data Overview")

    # Left Column: 4 Modalities
    add_card(slide2, 0.75, 1.4, 5.6, 5.3)
    col1_box = slide2.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(5.2), Inches(4.9))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Embedded Dataset Modalities"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)

    modalities = [
        ("Clinical Data", "Standardized scales (MoCA, Age, Demographics) autoencoded into a 32-dim latent space."),
        ("Genetic Markers", "PPMI genotyping arrays projected to a 32-dim latent space."),
        ("Structural MRI", "FreeSurfer regional brain volumes modeled as a synthetic KNN graph and passed through a 2-layer GCN to output a 64-dim embedding."),
        ("PET / DATScan", "Striatum dopamine transporter binding ratios (SBR) passed through a GAT (Graph Attention Network) to yield a 16-dim embedding.")
    ]
    for name, desc in modalities:
        p_item = tf1.add_paragraph()
        p_item.text = f"• {name}: "
        p_item.font.bold = True
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = TEXT_DARK
        p_item.font.name = "Arial"
        
        # Append description to the same paragraph
        run = p_item.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT_DARK
        p_item.space_after = Pt(10)

    # Right Column: Baseline Performance
    add_card(slide2, 6.75, 1.4, 5.8, 5.3)
    col2_box = slide2.shapes.add_textbox(Inches(6.95), Inches(1.6), Inches(5.4), Inches(4.9))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "14 Baseline Models Evaluated"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)

    p_desc = tf2.add_paragraph()
    p_desc.text = "We constructed a rigorous, leak-proof baseline (SMOTE + MICE imputation inside nested CV folds) evaluating 14 different classifiers and regressors:"
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = TEXT_DARK
    p_desc.space_after = Pt(10)

    models_list = [
        "Traditional ML: Random Forest, XGBoost, LightGBM, SVM, Ridge, AdaBoost, Naive Bayes, ElasticNet, KNN",
        "Deep Learning: Multi-Layer Perceptron (MLP), 1D-CNN, TabNet",
        "Graph & Federated GNNs: Patient similarity Graph Neural Network, and Simulated FedAvg GNN (3 clients)"
    ]
    for item in models_list:
        p_m = tf2.add_paragraph()
        p_m.text = f"✔ {item}"
        p_m.font.size = Pt(12)
        p_m.font.color.rgb = TEXT_DARK
        p_m.space_after = Pt(8)

    p_res = tf2.add_paragraph()
    p_res.text = "\nKey Baseline Results (Tabular Embeddings)"
    p_res.font.bold = True
    p_res.font.size = Pt(14)
    p_res.font.color.rgb = GOLD
    p_res.space_after = Pt(4)

    p_res_val = tf2.add_paragraph()
    p_res_val.text = "• Best Regression: XGBoost R² = 0.679, Random Forest CCC = 0.808\n• Best Classification: AdaBoost AUC = 0.561 (Tested on synthetic labels)\n• Total Patients/Rows in Intersection Dataset: 3,551 patients"
    p_res_val.font.size = Pt(12)
    p_res_val.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 3: Limitations of Baselines
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, LIGHT_GRAY)
    add_title(slide3, "Why Baseline Architectures are Insufficient for Top-Tier Journals")

    lims = [
        ("1. Simple Feature Concatenation", 
         "Concatenating embeddings assumes clinical, genetic, and imaging modalities operate independently. Reviewers expect a model that dynamically captures cross-talk and inter-modality biological correlations.",
         GOLD),
        ("2. Data Scale & Overfitting Risk", 
         "Training complex deep networks directly on classification/regression with only 3,551 medical samples results in high variance and overfitting. A robust similarity-based feature space is missing.",
         GOLD),
        ("3. Lack of Privacy & Hospital Simulation", 
         "PPMI is a multi-site clinical study. Standard central training ignores hospital-level data isolation. Top-tier medical informatics journals demand realistic Federated Learning evaluations.",
         GOLD)
    ]

    card_width = 3.64
    card_gap = 0.45
    left_start = 0.75
    for i, (title, desc, color) in enumerate(lims):
        c_left = left_start + i * (card_width + card_gap)
        add_card(slide3, c_left, 1.8, card_width, 4.6)
        
        tb = slide3.shapes.add_textbox(Inches(c_left + 0.15), Inches(2.0), Inches(card_width - 0.3), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Trebuchet MS"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(14)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_DARK
        p_d.space_after = Pt(12)

    # =========================================================================
    # SLIDE 4: Proposed CM-CADT Architecture
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, LIGHT_GRAY)
    add_title(slide4, "Proposed Architecture: Cross-Modal Co-Attentive Dual-Twin Network")

    # Left Column: Model Flow
    add_card(slide4, 0.75, 1.4, 5.6, 5.3)
    col1_box = slide4.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(5.2), Inches(4.9))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "How CM-CADT Operates"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)

    steps = [
        ("Step 1: Unified Projection", "Clinical, Genetic, GCN, and GAT outputs are projected into a unified 64-dimensional space as 'tokens'."),
        ("Step 2: Cross-Modal Attention", "A Multi-Head Self-Attention Transformer lets tokens attend to each other, computing complex inter-modality representations."),
        ("Step 3: Dual-Twin Encoders", "Weight-shared projection networks process patient pairs ($x_i, x_j$) into a contrastive latent space."),
        ("Step 4: Multi-Task Head", "The output is optimized jointly for Supervised Contrastive Loss, Diagnosis classification, and UPDRS-III regression.")
    ]
    for name, desc in steps:
        p_s = tf1.add_paragraph()
        p_s.text = f"✔ {name}: "
        p_s.font.bold = True
        p_s.font.size = Pt(12)
        p_s.font.color.rgb = TEXT_DARK
        
        run = p_s.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(12)
        p_s.space_after = Pt(8)

    # Right Column: The Tech Advantage
    add_card(slide4, 6.75, 1.4, 5.8, 5.3)
    col2_box = slide4.shapes.add_textbox(Inches(6.95), Inches(1.6), Inches(5.4), Inches(4.9))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Why This is a Journal-Grade Contribution"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)

    advantages = [
        ("Mathematical Novelty", "Instead of arbitrary concatenation, Multi-Head Attention allows clinical symptoms to dynamically query brain morphology tokens."),
        ("Effective Sample Size Boost", "By training on pairs (Dual-Twin) instead of individual patients, the effective training data expands quadratically ($O(N^2)$), mitigating overfitting."),
        ("Supervised Contrastive Loss (SupCon)", "Rather than unsupervised clustering, SupCon leverages diagnosis labels to push PD patients closer to each other while isolating Healthy Controls.")
    ]
    for name, desc in advantages:
        p_a = tf2.add_paragraph()
        p_a.text = f"★ {name}\n"
        p_a.font.bold = True
        p_a.font.size = Pt(13)
        p_a.font.color.rgb = GOLD
        
        run = p_a.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_DARK
        p_a.space_after = Pt(10)

    # =========================================================================
    # SLIDE 5: Site-Based Federated Learning
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, LIGHT_GRAY)
    add_title(slide5, "Pillar 2: Site-Based Federated Learning (FedAvg)")

    # Left box: Site isolation
    add_card(slide5, 0.75, 1.4, 5.6, 5.3)
    col1_box = slide5.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(5.2), Inches(4.9))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "Real-World Data Distribution"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)

    p_txt = tf1.add_paragraph()
    p_txt.text = "PPMI is naturally collected across dozens of global clinical sites. Rather than artificial random federated splits, we implement a realistic Site-Based Federated framework:\n\n" \
                 "• Clinic Sites as Local Clients: We group patients by their actual CLINICAL_SITE IDs in PPMI.\n" \
                 "• Non-IID Cohorts: Different hospitals have different scanner brands, genetic pools, and patient counts, introducing realistic non-IID challenges.\n" \
                 "• Multi-Center Evaluation: The model is validated on the ability to generalize to new, unseen hospitals."
    p_txt.font.size = Pt(13)
    p_txt.font.color.rgb = TEXT_DARK
    p_txt.space_after = Pt(12)

    # Right box: Federated loop
    add_card(slide5, 6.75, 1.4, 5.8, 5.3)
    col2_box = slide5.shapes.add_textbox(Inches(6.95), Inches(1.6), Inches(5.4), Inches(4.9))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Fed-CM-CADT Training Scheme"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.space_after = Pt(14)

    fed_steps = [
        ("1. Broadcast", "The central server broadcasts the initial CM-CADT model parameters to each hospital site client."),
        ("2. Local Optimization", "Each hospital client trains the model on its local patient cohort for a set number of epochs using SupCon and classification losses."),
        ("3. Parameter Aggregation", "Hospitals securely send only model updates/gradients (no patient data is shared) back to the central server."),
        ("4. FedAvg Update", "The server averages the client weights weighted by site sample sizes, updating the global model. Repeat for N rounds.")
    ]
    for step, desc in fed_steps:
        p_fs = tf2.add_paragraph()
        p_fs.text = f"• {step}: "
        p_fs.font.bold = True
        p_fs.font.size = Pt(12)
        p_fs.font.color.rgb = TEXT_DARK
        
        run = p_fs.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(12)
        p_fs.space_after = Pt(8)

    # =========================================================================
    # SLIDE 6: Explainable AI (XAI) Suite
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, LIGHT_GRAY)
    add_title(slide6, "Pillar 3: The 3-Level Explainable AI (XAI) Suite")

    xais = [
        ("Level 1: Cross-Modal Attention Maps", 
         "We extract the self-attention weights from the Transformer layer. This visualizes exactly how much attention the model places on structural brain regions (MRI/PET) during clinical and genetic queries.",
         TEAL),
        ("Level 2: Tabular SHAP Values", 
         "SHAP (Shapley Additive exPlanations) is applied to the Clinical and Genetic MLP encoders. This identifies individual input features (e.g. MoCA score, specific risk genes) driving the patient's latent projection.",
         TEAL),
        ("Level 3: Latent Proximity Biomarker", 
         "A patient's severity index is mathematically defined as their Euclidean distance to the centroid of the Healthy Control cluster in the dual-twin space. We map this distance to clinical UPDRS-III scores.",
         TEAL)
    ]

    card_width = 3.64
    card_gap = 0.45
    left_start = 0.75
    for i, (title, desc, color) in enumerate(xais):
        c_left = left_start + i * (card_width + card_gap)
        add_card(slide6, c_left, 1.8, card_width, 4.6)
        
        tb = slide6.shapes.add_textbox(Inches(c_left + 0.15), Inches(2.0), Inches(card_width - 0.3), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Trebuchet MS"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(14)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_DARK
        p_d.space_after = Pt(12)

    # =========================================================================
    # SLIDE 7: Publication Roadmap & Target Journals
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, NAVY)
    
    add_title(slide7, "Strategic Target Journals & Execution Roadmap", color=GOLD, top=0.5)

    # Left: Target Journals
    target_box = slide7.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
    tf_tj = target_box.text_frame
    tf_tj.word_wrap = True

    p = tf_tj.paragraphs[0]
    p.text = "Target Q1 Journals"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(14)

    journals = [
        ("IEEE Journal of Biomedical & Health Informatics (JBHI)", "Impact Factor: 7.7. Focus: multi-modal fusion, federated learning. Perfect match for Fed-CM-CADT."),
        ("Elsevier Computers in Biology and Medicine", "Impact Factor: 7.7. Focus: clinical application, comparisons, XAI. High interest in Parkinson's modeling."),
        ("IEEE Transactions on Medical Imaging (TMI)", "Impact Factor: 10.6. Very high image-processing bar. TVB or raw image features would be preferred here.")
    ]
    for j_name, j_desc in journals:
        p_j = tf_tj.add_paragraph()
        p_j.text = f"✔ {j_name}\n"
        p_j.font.bold = True
        p_j.font.size = Pt(12)
        p_j.font.color.rgb = TEXT_LIGHT
        
        run = p_j.add_run()
        run.text = j_desc
        run.font.bold = False
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(148, 163, 184)
        p_j.space_after = Pt(10)

    # Right: Checklist
    chk_box = slide7.shapes.add_textbox(Inches(6.75), Inches(1.5), Inches(5.8), Inches(5.0))
    tf_c = chk_box.text_frame
    tf_c.word_wrap = True

    p = tf_c.paragraphs[0]
    p.text = "Next Steps / Implementation Plan"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(14)

    roadmap_steps = [
        "1. Extract true diagnosis labels (APPRDX) and hospital IDs (CLINICAL_SITE) from PPMI raw files.",
        "2. Implement the CM-CADT model architecture and project all modalities to unified tokens in PyTorch.",
        "3. Setup the simulated Federated Learning loop grouping local data by site.",
        "4. Integrate SHAP, attention matrix logging, and UMAP clustering visualizer.",
        "5. Compile the comparison results table against the 14 baselines."
    ]
    for step in roadmap_steps:
        p_r = tf_c.add_paragraph()
        p_r.text = step
        p_r.font.size = Pt(12)
        p_r.font.color.rgb = TEXT_LIGHT
        p_r.space_after = Pt(8)

    # Save presentation
    os.makedirs("outputs/presentation", exist_ok=True)
    ppt_path = "outputs/presentation/Parkinson_DualTwin_Plan.pptx"
    prs.save(ppt_path)
    print(f"Presentation successfully saved to: {os.path.abspath(ppt_path)}")

if __name__ == "__main__":
    create_presentation()
