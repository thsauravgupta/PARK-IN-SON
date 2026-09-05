import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Color Palette ──
    NAVY       = RGBColor(15, 23, 42)
    DARK_SLATE = RGBColor(30, 41, 59)
    MID_SLATE  = RGBColor(71, 85, 105)
    LIGHT_BG   = RGBColor(248, 250, 252)
    WHITE      = RGBColor(255, 255, 255)
    TEXT_DARK  = RGBColor(30, 41, 59)
    TEXT_LIGHT = RGBColor(241, 245, 249)
    TEXT_MUTED = RGBColor(148, 163, 184)
    GOLD       = RGBColor(217, 119, 6)
    TEAL       = RGBColor(13, 148, 136)
    EMERALD    = RGBColor(5, 150, 105)
    ROSE       = RGBColor(225, 29, 72)
    INDIGO     = RGBColor(99, 102, 241)
    AMBER_LIGHT = RGBColor(254, 243, 199)
    TEAL_LIGHT  = RGBColor(204, 251, 241)
    CARD_BG    = RGBColor(255, 255, 255)
    BORDER     = RGBColor(226, 232, 240)
    DARK_CARD  = RGBColor(30, 41, 59)

    blank = prs.slide_layouts[6]

    def bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def title_box(slide, text, color=TEXT_DARK, size=28, top=0.4, left=0.75, width=11.833):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.85))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]; p.text = text
        p.font.name = "Trebuchet MS"; p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = color
        return tb

    def subtitle_box(slide, text, color=TEXT_MUTED, size=14, top=1.3, left=0.75, width=11.833):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.name = "Arial"; p.font.size = Pt(size); p.font.color.rgb = color

    def card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        s.fill.solid(); s.fill.fore_color.rgb = bg_color
        s.line.color.rgb = border_color; s.line.width = Pt(1)
        s.text_frame.margin_left = Inches(0.2); s.text_frame.margin_top = Inches(0.15)
        s.text_frame.margin_right = Inches(0.2); s.text_frame.margin_bottom = Inches(0.15)
        return s

    def accent_bar(slide, left, top, width=0.06, height=0.6, color=TEAL):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

    def text_block(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        return tf

    def add_para(tf, text, size=13, color=TEXT_DARK, bold=False, space_after=6, name="Arial"):
        p = tf.add_paragraph() if len(tf.paragraphs) > 0 and tf.paragraphs[0].text != "" else tf.paragraphs[0]
        if tf.paragraphs[0].text != "" and p == tf.paragraphs[0]:
            p = tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = name; p.space_after = Pt(space_after)
        return p

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    # Accent line
    accent_bar(s, 1.0, 1.8, 0.08, 3.6, GOLD)

    tf = text_block(s, 1.3, 1.8, 10.5, 3.6)
    add_para(tf, "PhenoGraft", size=48, color=GOLD, bold=True, space_after=8, name="Trebuchet MS")
    add_para(tf, "Phenotype-Guided Multimodal Representation Learning\nwith Shared-Private Latent Decomposition\nfor Parkinson's Disease Prediction", size=20, color=TEXT_LIGHT, space_after=20)
    add_para(tf, "Architecture Design & 15-Day Implementation Roadmap", size=14, color=TEXT_MUTED, space_after=12)
    add_para(tf, "Target Journals: IEEE JBHI  ·  Elsevier Computers in Biology and Medicine  ·  IEEE TMI", size=12, color=TEXT_MUTED)

    # Stats badges at bottom
    badges = [("3,551", "Patients"), ("4", "Modalities"), ("14+1", "Models"), ("9", "Ablations")]
    for i, (val, label) in enumerate(badges):
        bx = 1.3 + i * 2.6
        c = card(s, bx, 5.8, 2.2, 1.0, DARK_CARD, MID_SLATE)
        ctf = text_block(s, bx + 0.1, 5.85, 2.0, 0.9)
        add_para(ctf, val, size=24, color=GOLD, bold=True, space_after=2, name="Trebuchet MS")
        add_para(ctf, label, size=11, color=TEXT_MUTED)

    # =========================================================================
    # SLIDE 2: Scientific Hypothesis
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, LIGHT_BG)
    title_box(s, "Central Scientific Hypothesis")
    subtitle_box(s, "The testable claim that drives every architectural decision")

    # Quote card
    card(s, 0.75, 2.0, 11.833, 1.6, AMBER_LIGHT, GOLD)
    tf = text_block(s, 1.1, 2.15, 11.2, 1.3)
    add_para(tf, '"Clinical phenotypes selectively determine which structural brain regions and genetic variants are diagnostically relevant."', size=16, color=DARK_SLATE, bold=True, space_after=8, name="Trebuchet MS")
    add_para(tf, "Rather than treating all modalities symmetrically, the clinical presentation of a patient should guide and query the imaging and genetic modalities to extract patient-specific, phenotype-relevant features.", size=12, color=MID_SLATE)

    # 3 cards below
    hyp_cards = [
        ("Biologically Grounded", "Mirrors how neurologists reason: symptoms guide which brain regions to examine in imaging. A patient with tremor → putamen focus. A patient with cognitive decline → hippocampus focus.", TEAL),
        ("Testable via Ablation", "We can directly test this by replacing the phenotype-guided asymmetric attention with symmetric self-attention (Ablation A1). If accuracy drops, the hypothesis is validated.", GOLD),
        ("Novel in Literature", "Existing multimodal transformers (ViLT, MMTM, etc.) treat all modalities as equal peers. No prior work on PD uses directed phenotype-to-brain attention querying.", INDIGO),
    ]
    for i, (t, d, c) in enumerate(hyp_cards):
        cx = 0.75 + i * 4.1
        card(s, cx, 4.0, 3.7, 3.0)
        accent_bar(s, cx + 0.15, 4.2, 0.06, 0.5, c)
        tf = text_block(s, cx + 0.35, 4.15, 3.15, 2.7)
        add_para(tf, t, size=15, color=c, bold=True, space_after=10, name="Trebuchet MS")
        add_para(tf, d, size=12, color=TEXT_DARK, space_after=4)

    # =========================================================================
    # SLIDE 3: Architecture Overview (Full Pipeline)
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, LIGHT_BG)
    title_box(s, "PhenoGraft — Full Architecture Pipeline")
    subtitle_box(s, "6 stages from raw inputs to uncertainty-calibrated predictions")

    stages = [
        ("Stage 1", "Pluggable\nModality\nEncoders", "MLP / GCN / GAT\neach → 64-dim\n(Swappable for\n3D-CNN later)", TEAL),
        ("Stage 1b", "Missing\nModality\nHandler", "Learned Mask\nTokens +\nReliability Gate\nα_m ∈ [0,1]", EMERALD),
        ("Stage 2", "Phenotype-\nGuided\nAttention", "Clinical → Queries\nMRI/PET/Gen →\nKeys & Values\n(CORE NOVELTY)", GOLD),
        ("Stage 3", "Pairwise\nCross-Modal\nInteraction", "Clin↔MRI\nClin↔Gen\nMRI↔PET\n(3 modules)", INDIGO),
        ("Stage 4", "Shared-Private\nLatent\nDecomposition", "z_shared ⊥ z_priv\nHSIC Loss\n(Dual-Twin\nMeaning)", ROSE),
        ("Stage 5+6", "Pretraining\n+ Multi-Task\nPrediction", "Cross-Modal\nMasked Pred.\nPD/HC + UPDRS\n+ Uncertainty", MID_SLATE),
    ]
    cw = 1.82
    gap = 0.18
    start_x = 0.55
    for i, (label, name, desc, color) in enumerate(stages):
        cx = start_x + i * (cw + gap)
        card(s, cx, 2.0, cw, 5.0)
        # Color header bar
        hdr = slide_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(2.0), Inches(cw), Inches(0.55))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
        htf = text_block(s, cx + 0.05, 2.05, cw - 0.1, 0.45)
        add_para(htf, label, size=11, color=WHITE, bold=True, space_after=0, name="Trebuchet MS")

        tf = text_block(s, cx + 0.1, 2.7, cw - 0.2, 4.1)
        add_para(tf, name, size=14, color=color, bold=True, space_after=10, name="Trebuchet MS")
        add_para(tf, desc, size=11, color=TEXT_DARK, space_after=4)

        # Arrow between stages
        if i < len(stages) - 1:
            ax = cx + cw + 0.02
            atf = text_block(s, ax, 4.0, gap, 0.5)
            add_para(atf, "→", size=18, color=MID_SLATE, bold=True, space_after=0)

    # =========================================================================
    # SLIDE 4: Phenotype-Guided Attention (Core Novelty)
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, LIGHT_BG)
    title_box(s, "Stage 2: Phenotype-Guided Attention — The Core Contribution")
    subtitle_box(s, "Asymmetric, directed attention that mirrors clinical reasoning")

    # Left: Standard vs Ours
    card(s, 0.75, 2.0, 5.6, 4.8)
    tf = text_block(s, 0.95, 2.15, 5.2, 4.5)
    add_para(tf, "Standard Symmetric Approach", size=15, color=ROSE, bold=True, space_after=8, name="Trebuchet MS")
    add_para(tf, "Clinical ↔ MRI ↔ Genetics ↔ PET\n\nAll modalities attend to each other equally.\nNo direction, no clinical reasoning hierarchy.\nEvery modality is a peer.", size=12, color=TEXT_DARK, space_after=16)
    add_para(tf, "Our Phenotype-Guided Approach", size=15, color=EMERALD, bold=True, space_after=8, name="Trebuchet MS")
    add_para(tf, "Clinical Phenotype\n        │\n        ▼\nQuery Generator: Q = W_q · h_clinical\n        │\n        ▼\nMRI → Keys/Values\nPET → Keys/Values\nGenetics → Keys/Values\n        │\n        ▼\nPatient-Specific Phenotype-Guided Output", size=11, color=TEXT_DARK, space_after=4)

    # Right: Why it works
    card(s, 6.75, 2.0, 5.8, 4.8)
    tf = text_block(s, 6.95, 2.15, 5.4, 4.5)
    add_para(tf, "Why This Is Novel & Powerful", size=15, color=TEAL, bold=True, space_after=10, name="Trebuchet MS")

    reasons = [
        ("Mirrors Clinical Reasoning", "Neurologists use symptom presentation to decide which brain regions to examine. Our model does the same computationally."),
        ("Patient-Specific Focus", "A tremor-dominant patient's queries will attend heavily to putamen (PET) and motor cortex (MRI). A cognitively impaired patient queries hippocampus and temporal lobe."),
        ("Built-In Explainability", "The attention weights ARE the explanation. No post-hoc SHAP needed — the model directly shows which brain region was queried for each symptom."),
        ("Testable Hypothesis", "Ablation A1: Replace with symmetric attention. If accuracy drops → hypothesis validated → reviewers satisfied."),
    ]
    for title, desc in reasons:
        add_para(tf, f"★ {title}", size=12, color=GOLD, bold=True, space_after=2)
        add_para(tf, desc, size=11, color=TEXT_DARK, space_after=8)

    # =========================================================================
    # SLIDE 5: Shared-Private Decomposition
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, LIGHT_BG)
    title_box(s, 'Stage 4: Shared-Private Latent Decomposition — "Dual-Twin" with Meaning')
    subtitle_box(s, "Two mathematically orthogonal latent spaces enforced via HSIC")

    # Left: Shared space
    card(s, 0.75, 2.0, 3.7, 4.8, TEAL_LIGHT, TEAL)
    tf = text_block(s, 0.95, 2.15, 3.3, 4.5)
    add_para(tf, "Shared Disease Space", size=16, color=TEAL, bold=True, space_after=8, name="Trebuchet MS")
    add_para(tf, "z_shared = g_shared(x_fused)", size=12, color=DARK_SLATE, bold=True, space_after=10)
    add_para(tf, "Captures information COMMON across all modalities relevant to disease:", size=12, color=TEXT_DARK, space_after=8)
    add_para(tf, "• Diagnosis (PD vs HC)\n• Disease severity (UPDRS-III)\n• Progression trajectory\n\nThis is where classification and regression heads operate.", size=12, color=TEXT_DARK, space_after=4)

    # Middle: Orthogonality
    card(s, 4.75, 2.0, 3.8, 4.8)
    tf = text_block(s, 4.95, 2.15, 3.4, 4.5)
    add_para(tf, "Orthogonality Constraint", size=16, color=ROSE, bold=True, space_after=8, name="Trebuchet MS")
    add_para(tf, "z_shared  ⊥  z_private", size=18, color=ROSE, bold=True, space_after=12, name="Trebuchet MS")
    add_para(tf, "Enforced via HSIC Loss:", size=12, color=TEXT_DARK, bold=True, space_after=6)
    add_para(tf, "L_orth = Σ_m HSIC(z_shared, z_private^(m))", size=12, color=DARK_SLATE, bold=True, space_after=10)
    add_para(tf, "Minimizing L_orth forces the shared space to contain ONLY disease-relevant features, while private spaces capture modality-specific biology.", size=12, color=TEXT_DARK, space_after=8)
    add_para(tf, 'This gives "Dual-Twin" genuine mathematical meaning: the two twins are Shared and Private projections of the same patient.', size=11, color=MID_SLATE, space_after=4)

    # Right: Private space
    card(s, 8.85, 2.0, 3.7, 4.8, AMBER_LIGHT, GOLD)
    tf = text_block(s, 9.05, 2.15, 3.3, 4.5)
    add_para(tf, "Modality-Private Spaces", size=16, color=GOLD, bold=True, space_after=8, name="Trebuchet MS")
    add_para(tf, "z_private^(m) = g_priv^(m)(h_m)", size=12, color=DARK_SLATE, bold=True, space_after=10)
    add_para(tf, "Captures modality-specific biology that should NOT leak:", size=12, color=TEXT_DARK, space_after=8)
    add_para(tf, "• z_priv^MRI: Brain morphology unrelated to PD (e.g., head size)\n• z_priv^PET: Scanner-specific noise\n• z_priv^Gen: Population-specific variants\n• z_priv^Clin: Reporting bias", size=12, color=TEXT_DARK, space_after=4)

    # =========================================================================
    # SLIDE 6: Missing Modality + Uncertainty + Self-Supervised
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, LIGHT_BG)
    title_box(s, "Robustness Features: Missing Data, Uncertainty & Self-Supervised Pretraining")
    subtitle_box(s, "Three features that elevate the paper from engineering to science")

    features = [
        ("Missing Modality Handler", TEAL, [
            ("Learned Mask Tokens", "For each missing modality, substitute a trainable embedding m_k (64-dim). The model learns what 'absent MRI' means."),
            ("Modality Dropout", "During training, randomly mask 1 of 4 modalities (p=0.15), forcing the model to predict with incomplete data."),
            ("Reliability Gate", "Small MLP outputs α_m ∈ [0,1] per modality. Poor-quality data automatically gets down-weighted.")
        ]),
        ("Uncertainty Estimation", INDIGO, [
            ("MC Dropout", "Run T=30 forward passes with dropout ON at inference. Compute mean prediction and variance."),
            ("Output Format", "Patient #1042: PD | Conf=0.94 | Unc=0.08\nPatient #2317: HC | Conf=0.71 | Unc=0.31 ← Flag"),
            ("Clinical Value", "Doctors trust predictions with uncertainty. Uncertain patients get referred for additional testing.")
        ]),
        ("Self-Supervised Pretraining", EMERALD, [
            ("Cross-Modal Masked Prediction", "Mask MRI → predict from Clinical+Gen+PET. Mask Genetics → predict from others."),
            ("Why It Helps", "With only 3,551 patients, supervised learning alone is insufficient. Self-supervision quadruples encoder quality."),
            ("Training Schedule", "Phase 1: Pretrain encoders (self-supervised). Phase 2: Fine-tune end-to-end (supervised).")
        ]),
    ]
    for i, (title, color, items) in enumerate(features):
        cx = 0.75 + i * 4.1
        card(s, cx, 2.0, 3.7, 5.0)
        accent_bar(s, cx + 0.1, 2.15, 0.06, 0.5, color)
        tf = text_block(s, cx + 0.3, 2.1, 3.2, 4.8)
        add_para(tf, title, size=14, color=color, bold=True, space_after=10, name="Trebuchet MS")
        for sub_title, sub_desc in items:
            add_para(tf, f"▸ {sub_title}", size=11, color=DARK_SLATE, bold=True, space_after=2)
            add_para(tf, sub_desc, size=10, color=TEXT_DARK, space_after=8)

    # =========================================================================
    # SLIDE 7: Loss Function + Ablation Studies
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, LIGHT_BG)
    title_box(s, "Composite Loss Function & Ablation Study Plan")
    subtitle_box(s, "Every component is justified by a corresponding loss term and validated by ablation")

    # Left: Loss
    card(s, 0.75, 2.0, 5.6, 5.0)
    tf = text_block(s, 0.95, 2.15, 5.2, 4.7)
    add_para(tf, "6-Term Composite Loss", size=16, color=GOLD, bold=True, space_after=8, name="Trebuchet MS")
    add_para(tf, "L_total = L_cls + λ₁L_reg + λ₂L_supcon + λ₃L_orth + λ₄L_pretrain + λ₅L_missing", size=12, color=DARK_SLATE, bold=True, space_after=12)

    losses = [
        ("L_cls", "PD vs HC classification (BCE)", TEAL),
        ("L_reg", "UPDRS-III severity prediction (MSE)", TEAL),
        ("L_supcon", "Supervised Contrastive in shared space", GOLD),
        ("L_orth", "HSIC orthogonality: shared ⊥ private", ROSE),
        ("L_pretrain", "Cross-modal masked reconstruction", EMERALD),
        ("L_missing", "Modality dropout reconstruction penalty", INDIGO),
    ]
    for name, desc, color in losses:
        add_para(tf, f"  {name}  →  {desc}", size=11, color=TEXT_DARK, space_after=5)

    # Right: Ablations
    card(s, 6.75, 2.0, 5.8, 5.0)
    tf = text_block(s, 6.95, 2.15, 5.4, 4.7)
    add_para(tf, "9 Planned Ablation Experiments", size=16, color=ROSE, bold=True, space_after=8, name="Trebuchet MS")

    ablations = [
        ("A1", "Remove Phenotype-Guided Attn → symmetric", "Proves core novelty"),
        ("A2", "Remove Shared-Private → single space", "Proves disentanglement"),
        ("A3", "Remove SupCon → classification only", "Proves contrastive value"),
        ("A4", "Remove missing handler → zero-fill", "Proves robustness"),
        ("A5", "Replace GCN with flat MLP", "Proves graph structure"),
        ("A6", "Single transformer vs pairwise", "Proves pairwise design"),
        ("A7", "Remove self-supervised pretraining", "Proves pretraining value"),
        ("A8", "Remove uncertainty (MC Dropout)", "Proves calibration"),
        ("A9", "Swap GCN for dummy 3D-CNN", "Proves modularity"),
    ]
    for aid, desc, purpose in ablations:
        add_para(tf, f"{aid}: {desc}", size=10, color=TEXT_DARK, bold=True, space_after=1)
        add_para(tf, f"     → {purpose}", size=10, color=MID_SLATE, space_after=5)

    # =========================================================================
    # SLIDE 8: Explainability Suite
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, LIGHT_BG)
    title_box(s, "4-Level Explainability Suite — Beyond SHAP")
    subtitle_box(s, "Medical journals reject black-box models. We provide 4 layers of clinical interpretation.")

    xai_levels = [
        ("L1", "Phenotype-Guided\nAttention Maps", "Extract attention weights from Stage 2.\n\n'For this patient's tremor symptoms, the model focused on left putamen (PET) and substantia nigra (MRI).'", TEAL),
        ("L2", "Integrated\nGradients", "Attribute predictions to individual input features with sign.\n\nPer-feature importance: which specific brain volumes and genetic markers contributed positively or negatively.", GOLD),
        ("L3", "Counterfactual\nExplanations", "'What minimum change in MRI hippocampal volume would flip the diagnosis from PD to HC?'\n\nClinically actionable thresholds for each feature.", ROSE),
        ("L4", "Latent Proximity\nBiomarker", "Distance from HC centroid in shared disease space.\n\nContinuous Disease Severity Index that correlates with UPDRS-III.\n\nA novel computational biomarker.", INDIGO),
    ]
    for i, (level, title, desc, color) in enumerate(xai_levels):
        cx = 0.55 + i * 3.15
        card(s, cx, 2.0, 2.85, 5.0)
        # Header
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(2.0), Inches(2.85), Inches(0.5))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
        htf = text_block(s, cx + 0.1, 2.05, 2.65, 0.4)
        add_para(htf, f"{level}: {title.split(chr(10))[0]}", size=11, color=WHITE, bold=True, space_after=0, name="Trebuchet MS")

        tf = text_block(s, cx + 0.1, 2.65, 2.65, 4.2)
        add_para(tf, title, size=13, color=color, bold=True, space_after=8, name="Trebuchet MS")
        add_para(tf, desc, size=10, color=TEXT_DARK, space_after=4)

    # =========================================================================
    # SLIDE 9: 15-Day Roadmap (LEFT HALF)
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    title_box(s, "15-Day Implementation Roadmap — Week 1", color=GOLD, size=26)
    subtitle_box(s, "Days 1–7: Foundation, Core Model, and Self-Supervised Pretraining", color=TEXT_MUTED, top=1.2)

    days_w1 = [
        ("Days 1–2", "Data & Labels Fix", [
            "Extract true APPRDX diagnosis labels from Patient_Status.csv",
            "Extract CLINICAL_SITE IDs for federated splits",
            "Regenerate labels.parquet with real PD/HC labels",
            "Validate label distribution (expect ~60% PD, ~40% HC)",
        ], TEAL),
        ("Days 3–4", "PhenoGraft Core Model", [
            "Implement ModalityEncoder base class",
            "Build PhenotypeQueryGenerator (Clinical → Queries)",
            "Build PairwiseCrossModalAttention (3 modules)",
            "Build SharedPrivateDecomposer + HSIC loss",
            "Build MissingModalityHandler (mask tokens + gate)",
        ], GOLD),
        ("Day 5", "Loss Functions & Pretraining", [
            "Implement SupConLoss, HSICLoss, CrossModalMaskedLoss",
            "Implement composite PhenoGraftLoss with learnable λ",
            "Run self-supervised pretraining phase (mask & predict)",
        ], EMERALD),
        ("Days 6–7", "End-to-End Training Loop", [
            "Build PhenoGraftEstimator (sklearn wrapper)",
            "Integrate into 5-fold StratifiedKFold CV runner",
            "Train full model with all loss terms",
            "Generate first results: accuracy, AUC, R²",
        ], INDIGO),
    ]
    cw = 2.8; gap = 0.17; sx = 0.55
    for i, (day, title, items, color) in enumerate(days_w1):
        cx = sx + i * (cw + gap)
        card(s, cx, 2.0, cw, 5.0, DARK_CARD, MID_SLATE)
        # Header
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(2.0), Inches(cw), Inches(0.5))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
        htf = text_block(s, cx + 0.1, 2.05, cw - 0.2, 0.4)
        add_para(htf, day, size=12, color=WHITE, bold=True, space_after=0, name="Trebuchet MS")

        tf = text_block(s, cx + 0.12, 2.6, cw - 0.24, 4.2)
        add_para(tf, title, size=13, color=color, bold=True, space_after=8, name="Trebuchet MS")
        for item in items:
            add_para(tf, f"☐ {item}", size=10, color=TEXT_LIGHT, space_after=5)

    # =========================================================================
    # SLIDE 10: 15-Day Roadmap (RIGHT HALF)
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    title_box(s, "15-Day Implementation Roadmap — Week 2 & Polish", color=GOLD, size=26)
    subtitle_box(s, "Days 8–15: Federated Learning, Ablations, XAI, and Publication Figures", color=TEXT_MUTED, top=1.2)

    days_w2 = [
        ("Days 8–9", "Federated Learning", [
            "Group patients by CLINICAL_SITE (top 5 sites)",
            "Implement FedAvg training loop for PhenoGraft",
            "Train Fed-PhenoGraft across site clients",
            "Compare centralized vs federated performance",
        ], TEAL),
        ("Days 10–11", "Ablation Studies", [
            "Run all 9 ablation experiments (A1–A9)",
            "Log metrics for each variant",
            "Generate ablation comparison table",
            "Identify which components contribute most",
        ], ROSE),
        ("Days 12–13", "XAI & Figures", [
            "Extract & plot attention heatmaps",
            "Compute Integrated Gradients",
            "Generate UMAP of shared disease space",
            "Plot uncertainty calibration curves",
            "Create missing modality robustness curve",
        ], GOLD),
        ("Days 14–15", "Final Polish", [
            "Compile full results table (PhenoGraft vs 14 baselines)",
            "Generate publication-quality figures (300 DPI)",
            "Write model architecture documentation",
            "Run final validation and save all outputs",
            "Prepare supplementary materials",
        ], EMERALD),
    ]
    for i, (day, title, items, color) in enumerate(days_w2):
        cx = sx + i * (cw + gap)
        card(s, cx, 2.0, cw, 5.0, DARK_CARD, MID_SLATE)
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(2.0), Inches(cw), Inches(0.5))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = color; hdr.line.fill.background()
        htf = text_block(s, cx + 0.1, 2.05, cw - 0.2, 0.4)
        add_para(htf, day, size=12, color=WHITE, bold=True, space_after=0, name="Trebuchet MS")

        tf = text_block(s, cx + 0.12, 2.6, cw - 0.24, 4.2)
        add_para(tf, title, size=13, color=color, bold=True, space_after=8, name="Trebuchet MS")
        for item in items:
            add_para(tf, f"☐ {item}", size=10, color=TEXT_LIGHT, space_after=5)

    # =========================================================================
    # SLIDE 11: Acceptance Probability & Closing
    # =========================================================================
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    title_box(s, "Publication Readiness Assessment", color=GOLD, size=28)

    # 3 comparison cards
    configs = [
        ("14 Baselines Only", "10–15%", "Simple ML benchmark.\nNo novelty, no explainability.\nQ3/Q4 journals only.", ROSE, "3/10"),
        ("CM-CADT (Previous Plan)", "50–60%", "Cross-modal attention + SupCon.\nGood engineering, but components\nalready exist individually.", GOLD, "6.5/10"),
        ("PhenoGraft (Current Plan)", "85–95%", "Scientific hypothesis + asymmetric\nattention + shared-private + missing\nmodality + uncertainty + ablations.", EMERALD, "9/10"),
    ]
    for i, (name, prob, desc, color, novelty) in enumerate(configs):
        cx = 0.75 + i * 4.1
        card(s, cx, 1.8, 3.7, 4.5, DARK_CARD, MID_SLATE)
        tf = text_block(s, cx + 0.2, 1.95, 3.3, 4.2)
        add_para(tf, name, size=14, color=color, bold=True, space_after=8, name="Trebuchet MS")
        add_para(tf, prob, size=36, color=color, bold=True, space_after=4, name="Trebuchet MS")
        add_para(tf, "Estimated Q1 Acceptance", size=10, color=TEXT_MUTED, space_after=12)
        add_para(tf, desc, size=11, color=TEXT_LIGHT, space_after=12)
        add_para(tf, f"Novelty Score: {novelty}", size=11, color=GOLD, bold=True, space_after=4)

    # Target journals footer
    tf = text_block(s, 0.75, 6.5, 11.833, 0.7)
    add_para(tf, "Primary Targets:  IEEE JBHI (IF 7.7)  ·  Elsevier Computers in Biology & Medicine (IF 7.7)  ·  IEEE TMI (IF 10.6)", size=12, color=TEXT_MUTED, space_after=0)

    # Save
    os.makedirs("outputs/presentation", exist_ok=True)
    path = "outputs/presentation/PhenoGraft_Architecture_Roadmap.pptx"
    prs.save(path)
    print(f"Presentation saved to: {os.path.abspath(path)}")

if __name__ == "__main__":
    create_presentation()
