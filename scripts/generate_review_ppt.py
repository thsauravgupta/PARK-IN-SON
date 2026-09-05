import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # We will use a blank layout for all slides to have full control
    blank_layout = prs.slide_layouts[6]
    
    # Common function to add header/footer as requested by user
    def add_header_footer(slide, slide_num):
        # 21-07-2026
        tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(2), Inches(0.5))
        p1 = tb1.text_frame.paragraphs[0]
        p1.text = "21-07-2026"
        p1.font.size = Pt(12)
        p1.font.bold = True
        
        # SCOPE
        tb2 = slide.shapes.add_textbox(Inches(4.5), Inches(0.2), Inches(1), Inches(0.5))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = "SCOPE"
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER
        
        # Slide number (#)
        tb3 = slide.shapes.add_textbox(Inches(9.0), Inches(0.2), Inches(0.5), Inches(0.5))
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = str(slide_num)
        p3.font.size = Pt(12)
        p3.font.bold = True
        
    def add_title(slide, title_text):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        
    def add_content(slide, content_text, top=1.8, font_size=18):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content_text
        p.font.size = Pt(font_size)
        return tf

    # ==========================================
    # Slide 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    tb = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "<<Programme B.Tech/ M.Tech>>\n<<Course Code – Course Title>>"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nFed-PhenoGraft: Phenotype-Guided Asymmetric Cross-Modal Attention with Shared-Private Latent Decomposition for Federated Multi-Modal Parkinson's Disease Prediction\n"
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    p = tf.add_paragraph()
    p.text = "Team members:\n"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Saurav (<<Reg.No>>)\nAmit (<<Reg.No>>)\nShreeyam (<<Reg.No>>)\n"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "Faculty guide:\n"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "<<Guidename>>\n<To …>"
    p.font.size = Pt(16)

    # ==========================================
    # Slide 2: Aim
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide2, 2)
    add_title(slide2, "Aim")
    aim_text = (
        "To develop a federated, phenotype-guided multimodal representation learning framework "
        "(Fed-PhenoGraft) that improves the prediction of Parkinson's Disease (PD) diagnosis and severity. "
        "\n\nIt achieves this by asymmetrically querying imaging and genetic modalities based on clinical phenotypes, "
        "while simultaneously preserving data privacy across multiple clinical sites without centralizing raw patient data."
    )
    add_content(slide2, aim_text, font_size=20)

    # ==========================================
    # Slide 3: Abstract
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide3, 3)
    add_title(slide3, "Abstract")
    abs_text = (
        "Parkinson’s Disease (PD) diagnosis benefits significantly from multimodal data (clinical, genetics, MRI, PET). "
        "However, existing methods treat modalities symmetrically and require centralized data.\n\n"
        "We propose Fed-PhenoGraft, a novel federated learning framework featuring:\n"
        "1. Phenotype-Guided Asymmetric Attention: Clinical phenotypes query relevant structural and genetic features.\n"
        "2. Shared-Private Latent Decomposition: Uses HSIC orthogonality to disentangle disease-relevant info from modality noise.\n\n"
        "Evaluated on the PPMI cohort, Fed-PhenoGraft aims to achieve superior prognostic accuracy and patient-specific interpretable biomarkers, all while maintaining privacy via Federated Learning."
    )
    add_content(slide3, abs_text, font_size=18)

    # ==========================================
    # Slide 4: Literature Review
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide4, 4)
    add_title(slide4, "Literature Review")
    lit_text = (
        "• Modality Fusion: Combining heterogeneous data (MRI, PET, UPDRS, genetics) from the PPMI database consistently outperforms unimodal approaches for PD progression tracking [1].\n\n"
        "• Architectural Innovations: Recent models utilize cross-attention modules to align different data modalities and dynamically weight their importance, extracting more informative features [2].\n\n"
        "• Explainability (XAI): There is an increasing demand for explainable AI frameworks in PD diagnosis to identify affected brain regions and relevant genetic variants, moving away from \"black-box\" deep learning [3]."
    )
    add_content(slide4, lit_text, font_size=18)

    # ==========================================
    # Slide 5: Research Gap
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide5, 5)
    add_title(slide5, "Research Gap")
    gap_text = (
        "• Symmetric Treatment of Modalities: Current models treat all data sources as equal peers, lacking the biological motivation where clinical symptoms should actively guide the analysis of imaging data.\n\n"
        "• Entangled Latent Spaces: Existing fusion methods fail to cleanly separate disease-relevant information from modality-specific variations (e.g., scanner noise, demographic differences).\n\n"
        "• Data Privacy Issues: State-of-the-art multimodal PD models require centralizing sensitive multi-site patient records, which presents major legal and practical barriers.\n\n"
        "• Handling Missing Data: In real clinical settings, not all modalities are available; current models often struggle with incomplete patient records."
    )
    add_content(slide5, gap_text, font_size=18)

    # ==========================================
    # Slide 6: Objectives
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide6, 6)
    add_title(slide6, "Objectives")
    obj_text = (
        "1. To design a Phenotype-Guided Asymmetric Attention mechanism where clinical features selectively query MRI, PET, and genetics.\n\n"
        "2. To implement a Shared-Private Latent Decomposition using HSIC orthogonality to disentangle disease biomarkers from noise.\n\n"
        "3. To build a federated learning pipeline (FedAvg) that trains across multiple clinical sites without centralizing raw data.\n\n"
        "4. To incorporate robustness features (learned mask tokens for missing modalities, Monte Carlo Dropout for uncertainty).\n\n"
        "5. To provide a multi-level explainability suite (Attention Maps, Integrated Gradients, Counterfactuals)."
    )
    add_content(slide6, obj_text, font_size=18)

    # ==========================================
    # Slide 7: SDGs and Outcomes
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide7, 7)
    add_title(slide7, "Identified: Sustainable Development Goals & Outcomes")
    
    tf = add_content(slide7, "Identified: Sustainable Development Goals", font_size=20)
    p = tf.paragraphs[0]; p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• SDG 3: Good Health and Well-being (Improving diagnostic accuracy and personalized treatment for Parkinson's Disease)\n"
    p.text += "• SDG 9: Industry, Innovation, and Infrastructure (Developing advanced, privacy-preserving AI infrastructure for healthcare)\n"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Outcomes"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Journal (Scopus): Target publication in high-impact Q1 journals (e.g., IEEE JBHI, Elsevier Computers in Biology and Medicine).\n"
    p.text += "• Product/Others: A robust, federated AI framework applicable to other multimodal neurodegenerative datasets."
    p.font.size = Pt(18)

    # ==========================================
    # Slide 8: References
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide8, 8)
    add_title(slide8, "References")
    ref_text = (
        "[1] S. W. Akram and C. K, \"Enhancing Parkinson's Disease Staging: An Integrative Deep Learning Framework for Multimodal Feature Selection,\" Journal of Molecular Neuroscience, Mar 2026.\n\n"
        "[2] V. Awasthi et al., \"HyCoSwin-PD: An explainable hybrid ConvNeXtV2-Swin transformer framework for Parkinson's disease detection from neuroimaging,\" MethodsX, Jun 2026.\n\n"
        "[3] T. Zhi et al., \"MultimodalCNN-PD: a Parkinson's disease diagnostics framework using multimodal convolutional neural network,\" Frontiers in Aging Neuroscience, 2026."
    )
    add_content(slide8, ref_text, font_size=16)

    # Save presentation
    os.makedirs("outputs/presentation", exist_ok=True)
    path = "outputs/presentation/Project_Review_Presentation_v2.pptx"
    prs.save(path)
    print(f"Presentation saved to {os.path.abspath(path)}")

if __name__ == "__main__":
    create_presentation()
